# audit.py
"""
SAMO PPTX Audit Engine

Analyse un PPTX et produit audits/seance_<n>_audit.json avec :
 - le nombre de slides
 - les slides "spéciales" détectées (PLACEHOLDER / TABLE / TIMELINE)
 - un score d'intégrité (100 - pénalités)

Utilise le même extractor/classifier que converter_v2.py pour rester
cohérent entre l'audit et la conversion.
"""

import json
import sys
from pathlib import Path

from pptx import Presentation

from classifier import SlideType, classify_blocks
from extractor import extract_slide, find_boilerplate_texts

PENALTIES = {
    SlideType.GENERAL: 0,
    SlideType.PLACEHOLDER: 8,   # nécessite une action manuelle (ajouter un visuel)
    SlideType.TABLE: 2,
    SlideType.TIMELINE: 2,
}


def session_number(pptx_path):
    import re
    match = re.search(r"(\d+)", Path(pptx_path).stem)
    return match.group(1) if match else "1"


def audit_pptx(pptx_path):
    prs = Presentation(pptx_path)
    boilerplate = find_boilerplate_texts(prs)

    report = {
        "file": Path(pptx_path).name,
        "slides": len(prs.slides),
        "issues": [],
    }

    score = 100

    for index, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide(slide, boilerplate)
        slide_type = classify_blocks(slide_data["title"], slide_data["blocks"])

        if slide_type != SlideType.GENERAL:
            report["issues"].append({
                "slide": index,
                "type": slide_type,
                "title": slide_data["title"],
            })
            score -= PENALTIES.get(slide_type, 0)

    report["integrity_score"] = max(score, 0)
    return report


def main():
    pptx_file = sys.argv[1]
    report = audit_pptx(pptx_file)

    output_dir = Path("audits")
    output_dir.mkdir(exist_ok=True)

    n = session_number(pptx_file)
    output_file = output_dir / f"seance_{n}_audit.json"

    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    print(f"Audit généré : {output_file}")
    print(f"Score d'intégrité : {report['integrity_score']}/100")
    print(f"Slides analysées : {report['slides']}")
    for issue in report["issues"]:
        print(f"  - slide {issue['slide']:>2} [{issue['type']}] {issue['title']}")


if __name__ == "__main__":
    main()
