# extractor.py
"""
SAMO PPTX Extractor

Marche dans les shapes (y compris groupes et tableaux) d'une slide et
retourne le VRAI contenu texte, dans l'ordre. Filtre automatiquement le
"kicker" répété sur (presque) toutes les slides (fil d'ariane, nom du
formateur, version) en se basant sur la fréquence d'apparition, sans
rien coder en dur sur le contenu de la formation.
"""

import re
from collections import Counter, defaultdict

from pptx.enum.shapes import MSO_SHAPE_TYPE

NUMBER_STEP_RE = re.compile(r"^\d{2}$")
PLACEHOLDER_RE = re.compile(r"\[\s*screenshot|insérer ici|insert image", re.IGNORECASE)

BOILERPLATE_RATIO = 0.3  # un texte qui revient sur >30% des slides = kicker/footer
BOILERPLATE_PREFIX_LEN = 12  # regroupe les variantes d'un même kicker ("X" vs "X — FORMATEUR")


def _walk_shape(shape, blocks):
    """Ajoute à `blocks` les contenus texte/tableau trouvés dans `shape`."""

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            _walk_shape(sub_shape, blocks)
        return

    if getattr(shape, "has_table", False):
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if rows:
            blocks.append({"kind": "table", "rows": rows})
        return

    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            blocks.append({"kind": "text", "text": text})


def extract_slide_blocks(slide):
    """Retourne la liste ordonnée des blocs texte/tableau bruts d'une slide."""

    blocks = []
    for shape in slide.shapes:
        _walk_shape(shape, blocks)
    return blocks


def find_boilerplate_texts(presentation):
    """
    Détecte les textes qui reviennent sur une grosse proportion des slides
    (fil d'ariane, nom du formateur, tag de version...).

    Les variantes d'un même kicker (ex: "Chafik Moustaid" vs
    "Chafik Moustaid — Formateur") sont regroupées par préfixe commun avant
    de comparer au seuil, sinon chaque variante prise isolément peut
    sembler trop rare pour être détectée.
    """

    total_slides = len(presentation.slides)
    counts = Counter()

    for slide in presentation.slides:
        seen_this_slide = set()
        for block in extract_slide_blocks(slide):
            if block["kind"] == "text":
                seen_this_slide.add(block["text"])
        counts.update(seen_this_slide)

    groups = defaultdict(list)
    for text, count in counts.items():
        key = text[:BOILERPLATE_PREFIX_LEN].lower()
        groups[key].append((text, count))

    threshold = max(2, int(total_slides * BOILERPLATE_RATIO))

    boilerplate = set()
    for members in groups.values():
        group_total = sum(count for _, count in members)
        if group_total >= threshold:
            boilerplate.update(text for text, _ in members)

    return boilerplate


def extract_slide(slide, boilerplate):
    """
    Retourne un dict structuré pour une slide :
      {
        "title": str | None,
        "blocks": [ { "kind": "text"|"table", ... } ],   # sans le titre, sans le kicker
        "raw_blocks": [ ... ]                              # tout, non filtré
      }
    """

    raw_blocks = extract_slide_blocks(slide)

    body_blocks = [
        block for block in raw_blocks
        if not (block["kind"] == "text" and block["text"] in boilerplate)
    ]

    title = None
    remaining = body_blocks
    if remaining and remaining[0]["kind"] == "text":
        title = remaining[0]["text"]
        remaining = remaining[1:]

    return {
        "title": title,
        "blocks": remaining,
        "raw_blocks": raw_blocks,
    }
