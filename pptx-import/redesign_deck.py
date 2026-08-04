# redesign_deck.py
"""
SAMO PPTX Redesign — corrige le probleme "texte trop petit, case presque
vide" observe sur les diapositives generees (gabarit commun a tous les
cours : Support informatique, Word, Soutien Reseautique/Windows Server).

Principe : pour chaque zone de texte de CONTENU (on exclut le bandeau
d'en-tete/pied de page qui doit rester discret), on cherche la plus
grande taille de police qui remplit mieux la case SANS jamais deborder.

Historique des corrections :
  v2 - la v1 comptait 1 ligne par paragraphe (faux des qu'un paragraphe
       est assez long pour retourner a la ligne) -> chevauchements de
       texte. Corrige avec une estimation du nombre de lignes basee sur
       la largeur reelle de la case.
  v3 - trois bugs restants observes en production :
       1. Des cases "soeurs" de memes dimensions (ex. 4 cartes ENTREES/
          TRAITEMENT/MEMOIRE/SORTIES cote a cote) recevaient chacune une
          taille differente selon la longueur de LEUR texte -> incoherence
          visuelle flagrante entre cases pourtant identiques. Corrige en
          regroupant les formes de memes dimensions sur une meme slide et
          en appliquant a tout le groupe la taille minimale (sure) du
          groupe.
       2. Un mot seul plus large que la case (ex. une etiquette
          "Frequence" dans une case etroite) etait agrandi jusqu'a ne
          plus tenir sur une ligne, forcant une coupure AU MILIEU du mot
          ("Frequenc" / "e"). Corrige en imposant que le mot le plus long
          d'un texte tienne toujours sur une seule ligne a la taille
          choisie.
       3. De petites icones/pictos (fleches "▶", cases a cocher "☐",
          etoiles "★"...) etaient traitees comme du texte de contenu et
          agrandies, ce qui les faisait paraitre mal alignees avec leurs
          elements decoratifs voisins. Corrige en detectant et excluant
          les formes dont le texte est un simple symbole/picto.
  v4 - chevauchement titre/corps observe (ex. carte "Eliminer le Superflu
       (Bloatware)" dont le titre en gras se retrouvait recouvert par le
       paragraphe du dessous). Cause : le gabarit utilise deliberement
       une case de "corps" TRES haute (bien plus haute que necessaire)
       positionnee EN DESSOUS/DERRIERE une case de "titre" plus petite,
       le tout ancre au centre (MIDDLE). Avec un texte court, le centrage
       vertical dans une case surdimensionnee suffisait a rester sous le
       titre a la taille d'origine, mais agrandir la police (le but meme
       du redesign) fait grossir le bloc de texte centre et remonte son
       sommet jusqu'a chevaucher le titre. Corrige en detectant les
       paires titre/corps qui se chevauchent geometriquement : la case de
       corps passe en ancrage HAUT avec une marge superieure forcee sous
       le titre, et son calcul de taille de police utilise desormais la
       hauteur reellement disponible sous le titre (pas la hauteur totale
       de la case) pour ne jamais deborder.

Usage :
  python redesign_deck.py chemin/vers/Seance_N.pptx
  -> modifie le fichier sur place (fonctionne uniquement si le fichier
     est sur un disque local, ex. /tmp — pas sur un dossier reseau/sync)
"""

import math
import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt

# Bandeau haut (logo + titre de seance + formateur + version) et bandeau
# bas (fil d'ariane) : on ne touche pas a ces zones, elles doivent rester
# discretes. Coordonnees stables sur tout le gabarit (12188825 x 6858000).
HEADER_LIMIT_EMU = 620000
FOOTER_START_EMU = 6650000

# Largeur moyenne d'un caractere, en fraction de la taille de police
# (approximation raisonnable pour une police proportionnelle type Arial/
# Calibri, texte francais avec majuscules/minuscules melangees).
AVG_CHAR_WIDTH_FACTOR = 0.52
LINE_HEIGHT_FACTOR = 1.25

# Une forme dont le texte (tous paragraphes confondus) ne contient aucune
# lettre/chiffre est consideree comme un picto/icone decoratif (fleches,
# coches, etoiles...) et n'est jamais redimensionnee.
_HAS_ALNUM_RE = re.compile(r"[0-9A-Za-zÀ-ÖØ-öø-ÿ]")


def is_chrome(shape):
    if shape.top is None:
        return True
    if shape.top < HEADER_LIMIT_EMU:
        return True
    if shape.top > FOOTER_START_EMU:
        return True
    return False


def is_icon_glyph(texts):
    joined = "".join(texts)
    if len(joined) <= 2 and not _HAS_ALNUM_RE.search(joined):
        return True
    return False


def paragraph_texts(text_frame):
    texts = []
    for p in text_frame.paragraphs:
        t = "".join(r.text for r in p.runs).strip()
        if t:
            texts.append(t)
    return texts


def first_font_size_pt(text_frame):
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip() and r.font.size:
                return r.font.size.pt
    return None


def longest_word_len(texts):
    longest = 0
    for t in texts:
        for word in t.split():
            longest = max(longest, len(word))
    return longest


def is_bold_text(text_frame):
    for p in text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip() and r.font.bold:
                return True
    return False


def char_width_factor(bold):
    # Le gras est sensiblement plus large ; sous-estimer ca cause des
    # coupures de mot au milieu (observe en production).
    return 0.60 if bold else AVG_CHAR_WIDTH_FACTOR


def usable_box_dims(shape, top_override_emu=None):
    """Largeur/hauteur utilisables en pt, en retirant les marges internes
    du cadre de texte (91440 EMU par defaut de chaque cote sinon) -
    les ignorer causait des coupures de mot au milieu (le texte debordait
    en realite alors que le calcul le pensait tout juste tenu).

    Si top_override_emu est fourni (cas d'une case "corps" recouverte par
    un "titre" au-dessus), la hauteur utilisable est calculee a partir de
    cette position de depart plus basse, pas du sommet reel de la case."""
    tf = shape.text_frame
    margin_l = tf.margin_left if tf.margin_left is not None else 91440
    margin_r = tf.margin_right if tf.margin_right is not None else 91440
    margin_t = tf.margin_top if tf.margin_top is not None else 45720
    margin_b = tf.margin_bottom if tf.margin_bottom is not None else 45720
    usable_w = (shape.width - margin_l - margin_r) / 12700
    if top_override_emu is not None and top_override_emu > shape.top:
        effective_height = (shape.top + shape.height) - top_override_emu
        usable_h = (effective_height - margin_b) / 12700
    else:
        usable_h = (shape.height - margin_t - margin_b) / 12700
    return max(usable_w, 1), max(usable_h, 1)


def bbox_overlap_width(a_left, a_width, b_left, b_width):
    lo = max(a_left, b_left)
    hi = min(a_left + a_width, b_left + b_width)
    return max(hi - lo, 0)


def detect_header_pairs(shapes):
    """Detecte les paires "titre au-dessus / corps surdimensionne en
    dessous" qui se chevauchent geometriquement dans le gabarit (le corps
    est une case bien plus haute que necessaire, positionnee derriere le
    titre). Retourne un dict {id(shape body): header_bottom_emu} pour
    chaque case de corps concernee."""
    pairs = {}
    for b in shapes:
        if b.top is None or b.height is None:
            continue
        best_header_bottom = None
        for a in shapes:
            if a is b or a.top is None or a.height is None:
                continue
            if a.height >= b.height * 0.5:
                continue  # pas assez "petit" pour etre un titre
            # Le titre doit demarrer dans le premier tiers du corps.
            if not (b.top - 22860 <= a.top <= b.top + b.height * 0.3):
                continue
            overlap_h = max(0, min(a.top + a.height, b.top + b.height) - max(a.top, b.top))
            if overlap_h <= 0:
                continue
            overlap_w = bbox_overlap_width(a.left, a.width, b.left, b.width)
            if overlap_w < 0.5 * min(a.width, b.width):
                continue
            header_bottom = a.top + a.height
            if best_header_bottom is None or header_bottom > best_header_bottom:
                best_header_bottom = header_bottom
        if best_header_bottom is not None:
            pairs[b.shape_id] = best_header_bottom
    return pairs


def detect_horizontal_overlaps(shapes):
    """Detecte les formes dont la case deborde deja, par construction du
    gabarit, sur une case voisine placee a sa droite (ex. une etiquette
    "PARTITIONS MAX EN GPT" dont la case chevauche la carte a cote). Ce
    defaut peut exister meme sans agrandissement de police (largeur de
    case mal calibree a l'origine). Retourne {shape_id: safe_right_emu}."""
    overlaps = {}
    for b in shapes:
        if b.left is None or b.width is None:
            continue
        b_right = b.left + b.width
        best_safe_right = None
        for a in shapes:
            if a is b or a.left is None:
                continue
            if a.left <= b.left:
                continue  # on ne traite que les voisins strictement a droite
            v_overlap = max(0, min(b.top + b.height, a.top + a.height) - max(b.top, a.top))
            if v_overlap <= 0:
                continue
            if b_right <= a.left:
                continue  # pas de veritable intrusion
            if best_safe_right is None or a.left < best_safe_right:
                best_safe_right = a.left
        if best_safe_right is not None:
            overlaps[b.shape_id] = best_safe_right
    return overlaps


def shrink_to_fit(texts, box_width_pt, box_height_pt, start_pt, cwf, min_pt=8.0):
    """Reduit la taille de police (meme en dessous de la taille d'origine
    si necessaire) jusqu'a ce que le texte tienne dans une largeur/hauteur
    de case reduite (cas d'un chevauchement deja present dans le gabarit
    d'origine, independant de tout agrandissement)."""
    size = start_pt
    step = 0.5
    while size > min_pt:
        if word_fits(texts, box_width_pt, size, cwf):
            lines = estimated_lines(texts, box_width_pt, size, cwf)
            needed_pt = size * LINE_HEIGHT_FACTOR * lines
            if needed_pt <= box_height_pt:
                return round(size, 1)
        size -= step
    return min_pt


def estimated_lines(texts, box_width_pt, font_pt, cwf):
    chars_per_line = max(int(box_width_pt / (font_pt * cwf)), 1)
    total = 0
    for t in texts:
        total += max(1, math.ceil(len(t) / chars_per_line))
    return total


def word_fits(texts, box_width_pt, font_pt, cwf):
    # Le mot le plus long doit tenir sur une seule ligne a cette taille,
    # sinon le moteur de rendu le coupera au milieu (tres moche).
    longest = longest_word_len(texts)
    if longest == 0:
        return True
    return longest * font_pt * cwf <= box_width_pt * 0.96


def fits(texts, box_width_pt, box_height_pt, font_pt, target_occupancy, cwf):
    if not word_fits(texts, box_width_pt, font_pt, cwf):
        return False
    lines = estimated_lines(texts, box_width_pt, font_pt, cwf)
    needed_pt = font_pt * LINE_HEIGHT_FACTOR * lines
    return needed_pt <= target_occupancy * box_height_pt


def best_font_size(texts, box_width_pt, box_height_pt, old_pt, target_occupancy, max_factor, cwf):
    # Recherche la plus grande taille (par pas de 0.5pt) qui tient dans
    # la case, entre old_pt et old_pt * max_factor.
    hi = old_pt * max_factor
    best = old_pt
    size = old_pt
    step = 0.5
    while size <= hi:
        if fits(texts, box_width_pt, box_height_pt, size, target_occupancy, cwf):
            best = size
        else:
            break
        size += step
    return round(best, 1)


def current_occupancy(texts, box_width_pt, box_height_pt, font_pt, cwf):
    lines = estimated_lines(texts, box_width_pt, font_pt, cwf)
    needed_pt = font_pt * LINE_HEIGHT_FACTOR * lines
    if box_height_pt <= 0:
        return 1.0
    return needed_pt / box_height_pt


def target_for_occupancy(occ):
    if occ < 0.30:
        return 0.62, 2.0
    elif occ < 0.55:
        return 0.72, 1.5
    else:
        return 0.82, 1.2


def compute_candidate_size(shape, header_bottom_emu=None):
    """Retourne (old_pt, new_pt, texts) pour une forme de texte, sans rien
    modifier. None si la forme doit etre ignoree (pas de texte/police, ou
    icone decorative).

    header_bottom_emu : si cette case de corps est recouverte par un
    titre au-dessus (cf. detect_header_pairs), la hauteur utilisable pour
    le calcul de taille est limitee a l'espace reellement disponible sous
    ce titre, pour ne jamais faire deborder le texte dedans."""
    tf = shape.text_frame
    old_pt = first_font_size_pt(tf)
    if old_pt is None:
        return None

    texts = paragraph_texts(tf)
    if not texts or is_icon_glyph(texts):
        return None

    box_width_pt, box_height_pt = usable_box_dims(shape, top_override_emu=header_bottom_emu)
    cwf = char_width_factor(is_bold_text(tf))

    occ = current_occupancy(texts, box_width_pt, box_height_pt, old_pt, cwf)
    target, max_factor = target_for_occupancy(occ)
    new_pt = best_font_size(texts, box_width_pt, box_height_pt, old_pt, target, max_factor, cwf)
    new_pt = max(new_pt, old_pt)

    return old_pt, new_pt, texts


def apply_font_size(shape, new_pt, header_bottom_emu=None):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            if r.text.strip() and r.font.size:
                r.font.size = Pt(new_pt)

    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass

    if header_bottom_emu is not None and header_bottom_emu > shape.top:
        # Case de corps recouverte par un titre : ancrage HAUT + marge
        # superieure poussee sous le titre, pour ne jamais chevaucher.
        try:
            tf.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass
        gap = 45720  # ~0.05" d'espace visuel entre le titre et le corps
        tf.margin_top = (header_bottom_emu - shape.top) + gap
    else:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass


def redesign_table(shape):
    table = shape.table
    col_widths_pt = [col.width / 12700 for col in table.columns]

    for row in table.rows:
        for col_idx, cell in enumerate(row.cells):
            tf = cell.text_frame
            old_pt = first_font_size_pt(tf)
            if old_pt is None:
                continue
            texts = paragraph_texts(tf)
            if texts and is_icon_glyph(texts):
                continue
            # Boost modeste et sans risque pour les tableaux (colonnes
            # etroites -> le retour a la ligne y est deja plus sensible).
            new_pt = round(old_pt * 1.15, 1)
            cwf = char_width_factor(is_bold_text(tf))
            margin_l = tf.margin_left if tf.margin_left is not None else 91440
            margin_r = tf.margin_right if tf.margin_right is not None else 91440
            col_width_emu = None
            if col_idx < len(col_widths_pt):
                col_width_emu = col_widths_pt[col_idx] * 12700
            box_width_pt = (
                max((col_width_emu - margin_l - margin_r), 12700) / 12700
                if col_width_emu
                else None
            )
            if texts and box_width_pt and not word_fits(texts, box_width_pt, new_pt, cwf):
                new_pt = old_pt
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.text.strip() and r.font.size:
                        r.font.size = Pt(new_pt)


def redesign_slide(slide):
    text_shapes = [
        s for s in slide.shapes
        if not is_chrome(s) and s.has_text_frame and s.text_frame.text.strip()
    ]
    header_pairs = detect_header_pairs(text_shapes)
    h_overlaps = detect_horizontal_overlaps(text_shapes)

    candidates = []  # list of (shape, old_pt, new_pt, texts, header_bottom_emu)
    for shape in slide.shapes:
        if is_chrome(shape):
            continue
        if shape.has_table:
            redesign_table(shape)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            header_bottom_emu = header_pairs.get(shape.shape_id)
            result = compute_candidate_size(shape, header_bottom_emu=header_bottom_emu)
            if result is not None:
                old_pt, new_pt, texts = result
                safe_right_emu = h_overlaps.get(shape.shape_id)
                if safe_right_emu is not None:
                    # Chevauchement deja present sur une case voisine (avec
                    # ou sans agrandissement) : on reduit la police, meme
                    # sous la taille d'origine, pour l'eliminer.
                    tf = shape.text_frame
                    margin_l = tf.margin_left if tf.margin_left is not None else 91440
                    margin_r = tf.margin_right if tf.margin_right is not None else 91440
                    safe_width_pt = max((safe_right_emu - shape.left - margin_l - margin_r) - 45720, 12700) / 12700
                    _, box_height_pt = usable_box_dims(shape, top_override_emu=header_bottom_emu)
                    cwf = char_width_factor(is_bold_text(tf))
                    new_pt = shrink_to_fit(texts, safe_width_pt, box_height_pt, new_pt, cwf)
                candidates.append((shape, old_pt, new_pt, texts, header_bottom_emu))

    # Regroupe les formes "soeurs" (memes dimensions exactes -> cartes
    # repetees en ligne/grille) pour leur appliquer une taille UNIFORME
    # (la plus petite du groupe, donc toujours sure) plutot qu'une taille
    # individuelle qui romprait la coherence visuelle entre cartes. Les
    # cases "corps recouvert par un titre" ne sont jamais regroupees avec
    # d'autres (chacune a son propre espace disponible sous son titre).
    groups = {}
    for idx, (shape, old_pt, new_pt, texts, header_bottom_emu) in enumerate(candidates):
        if header_bottom_emu is not None or shape.shape_id in h_overlaps:
            continue
        key = (shape.width, shape.height)
        groups.setdefault(key, []).append(idx)

    for idx, (shape, old_pt, new_pt, texts, header_bottom_emu) in enumerate(candidates):
        key = (shape.width, shape.height)
        sibling_indices = groups.get(key, [])
        if header_bottom_emu is None and shape.shape_id not in h_overlaps and len(sibling_indices) > 1:
            uniform_pt = min(candidates[i][2] for i in sibling_indices)
            apply_font_size(shape, uniform_pt)
        else:
            apply_font_size(shape, new_pt, header_bottom_emu=header_bottom_emu)


def redesign(pptx_path):
    pptx_path = Path(pptx_path)
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        redesign_slide(slide)

    backup = pptx_path.with_suffix(pptx_path.suffix + ".bak")
    if not backup.exists():
        shutil.copy(pptx_path, backup)

    prs.save(pptx_path)
    return pptx_path


def main():
    pptx_file = sys.argv[1]
    out = redesign(pptx_file)
    print(f"Redesign applique : {out}")


if __name__ == "__main__":
    main()
